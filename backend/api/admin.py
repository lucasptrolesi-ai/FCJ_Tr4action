import io

from typing import List

from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Depends
from pydantic import BaseModel
from pptx import Presentation

from core.rag_engine import rag_engine, KnowledgeDoc
from core.config import METADATA_PATH
from core.auth import require_admin

router = APIRouter(prefix="/admin", tags=["admin"])


class KnowledgeStats(BaseModel):
    docs: int
    steps: List[str]


@router.get("/knowledge", response_model=KnowledgeStats)
def get_knowledge_stats(admin: dict = Depends(require_admin)) -> KnowledgeStats:
    """Protegido: apenas criadores de conteúdo FCJ"""
    stats = rag_engine.get_stats()
    return KnowledgeStats(**stats)


@router.post("/upload-pptx")
async def upload_pptx(
    step: str = Form(...),
    files: List[UploadFile] = File(...),
    admin: dict = Depends(require_admin),
):
    """
    Protegido: Recebe um ou mais PPTX da trilha, extrai textos dos slides
    e adiciona na base de conhecimento (RAG).
    Apenas criadores de conteúdo FCJ podem fazer upload.
    """
    docs: List[KnowledgeDoc] = []

    for f in files:
        if not f.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Apenas arquivos .pptx são permitidos.")

        content = await f.read()
        prs = Presentation(io.BytesIO(content))  # type: ignore

        slide_texts = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text:
                    texts.append(shp.text)
            if texts:
                joined = " ".join(texts).strip()
                slide_texts.append(f"Slide {idx}: {joined}")

        if not slide_texts:
            continue

        full_text = "\n".join(slide_texts)
        doc = KnowledgeDoc(
            id=f.filename,
            step=step,
            title=f.filename,
            text=full_text,
        )
        docs.append(doc)

    if not docs:
        raise HTTPException(status_code=400, detail="Não foi possível extrair texto dos PPTX enviados.")

    rag_engine.add_documents(docs)

    stats = rag_engine.get_stats()
    return {
        "status": "ok",
        "added": len(docs),
        "docs_total": stats["docs"],
        "steps": stats["steps"],
    }


@router.post("/reload")
def reload_knowledge(admin: dict = Depends(require_admin)):
    """Protegido: apenas criadores de conteúdo FCJ"""
    rag_engine.reload()
    stats = rag_engine.get_stats()
    return {"status": "ok", "message": "Base recarregada.", "stats": stats}
